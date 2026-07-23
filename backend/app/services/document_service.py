import os
import uuid
from pathlib import Path
from sqlalchemy.orm import Session

from app.core.config import settings
from app.models.document import Document, DocumentChunk
from app.models.enums import DocumentStatus, DocumentType
from app.ai.embeddings.embedder import get_embedder


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    if len(text.strip()) < 20:
        # Likely a scanned PDF -- fall back to OCR
        text = _ocr_pdf(path)
    return text


def _ocr_pdf(path: str) -> str:
    from pdf2image import convert_from_path
    import pytesseract

    images = convert_from_path(path)
    return "\n".join(pytesseract.image_to_string(img) for img in images)


def _extract_docx(path: str) -> str:
    import docx

    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)


def _extract_txt_or_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


EXTRACTORS = {
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    "text/plain": _extract_txt_or_md,
    "text/markdown": _extract_txt_or_md,
}


def clean_text(text: str) -> str:
    lines = [ln.strip() for ln in text.splitlines()]
    lines = [ln for ln in lines if ln]
    return "\n".join(lines)


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> list[str]:
    """Simple sliding-window chunker over characters. Swap for a
    token-aware / semantic chunker (e.g. LangChain's RecursiveCharacterTextSplitter)
    for production-grade quality."""
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunks.append(text[start:end])
        if end == n:
            break
        start = end - overlap
    return [c.strip() for c in chunks if c.strip()]


def process_document(db: Session, document: Document) -> None:
    """Runs the full pipeline: extract -> clean -> chunk -> embed -> store.
    Called synchronously after upload here for simplicity; in production this
    should be pushed to a background worker (Celery/RQ) so uploads don't
    block the HTTP request."""
    try:
        extractor = EXTRACTORS.get(document.mime_type, _extract_txt_or_md)
        raw_text = extractor(document.file_path)
        cleaned = clean_text(raw_text)
        chunks = chunk_text(cleaned)

        embedder = get_embedder()
        vectors = embedder.embed(chunks) if chunks else []

        for idx, (chunk, vector) in enumerate(zip(chunks, vectors)):
            db.add(
                DocumentChunk(
                    id=uuid.uuid4(),
                    company_id=document.company_id,
                    document_id=document.id,
                    department_id=document.department_id,
                    chunk_index=idx,
                    content=chunk,
                    embedding=vector,
                    doc_type=document.doc_type,
                    meta={"version": document.version, "tags": document.tags or []},
                )
            )
        document.status = DocumentStatus.READY
        db.commit()
    except Exception as e:  # noqa: BLE001
        document.status = DocumentStatus.FAILED
        db.commit()
        raise e


def save_upload(company_id: uuid.UUID, filename: str, content: bytes) -> str:
    company_dir = Path(settings.UPLOAD_DIR) / str(company_id)
    company_dir.mkdir(parents=True, exist_ok=True)
    safe_name = f"{uuid.uuid4()}_{filename}"
    file_path = company_dir / safe_name
    file_path.write_bytes(content)
    return str(file_path)
